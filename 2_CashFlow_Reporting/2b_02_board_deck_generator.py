!pip install python-pptx 

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd

metrics = {
    'Month': 'November 2025',
    'Revenue': 150000,
    'Net_Income': 25000,
    'Cash_Balance': 85000,
    'Runway_Months': 4.5
}

prs = Presentation()

slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Monthly Financial Review"
subtitle.text = f"Period: {metrics['Month']}\nPrepared by: Outsourced CFO Office"

slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Executive Dashboard"

left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

p = tf.add_paragraph()
p.text = f"Revenue: ${metrics['Revenue']:,.0f}"
p.font.size = Pt(28)
p.font.bold = True

p = tf.add_paragraph()
p.text = f"Net Income: ${metrics['Net_Income']:,.0f}"
p.font.size = Pt(28)
if metrics['Net_Income'] > 0:
    p.font.color.rgb = RGBColor(0, 128, 0) 
else:
    p.font.color.rgb = RGBColor(255, 0, 0) 

p = tf.add_paragraph()
p.text = f"Cash on Hand: ${metrics['Cash_Balance']:,.0f} ({metrics['Runway_Months']} Months Runway)"
p.font.size = Pt(28)

output_name = 'Monthly_Board_Deck.pptx'
prs.save(output_name)
print(f"Presentation saved as {output_name}")
